from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import List, Dict, Any, Optional
import os
import re
from app.constants.constants import PPTXConstants, FilePaths, SlideLayoutType, PresentationTheme


class PPTXCreator:
    def __init__(self):
        self.presentation = None
        self.config = None
    
    def create_presentation(self, slides_data: List[Dict[str, Any]], topic: str, config: Optional[Dict[str, Any]] = None) -> str:
        """Create a PowerPoint presentation from slide data with configuration options"""
        
        # Store config for use throughout the presentation creation
        self.config = config or {}
        
        # Create new presentation
        self.presentation = Presentation()
        
        # Set slide size based on aspect ratio (default to 16:9)
        aspect_ratio = self.config.get('aspect_ratio', PPTXConstants.DEFAULT_ASPECT_RATIO)
        if aspect_ratio == '16:9':
            self.presentation.slide_width = Inches(13.33)
            self.presentation.slide_height = Inches(7.5)
        elif aspect_ratio == '4:3':
            self.presentation.slide_width = Inches(10)
            self.presentation.slide_height = Inches(7.5)
            
        # Process each slide according to its type
        for slide_data in slides_data:
            self._create_slide(slide_data)
        
        # Save presentation
        sanitized_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_').lower()
        filename = f"{sanitized_topic}.pptx"
        filepath = os.path.join(FilePaths.PRESENTATIONS_DIR, filename)
        
        # Ensure presentations directory exists
        os.makedirs(FilePaths.PRESENTATIONS_DIR, exist_ok=True)
        
        self.presentation.save(filepath)
        return filepath
    
    def _create_slide(self, slide_data: Dict[str, Any]):
        """Create a single slide based on slide data and type"""
        
        slide_type = slide_data.get("slide_type", SlideLayoutType.BULLET_POINTS.value)
        
        if slide_type == SlideLayoutType.TITLE.value:
            self._create_title_slide(slide_data)
        elif slide_type == SlideLayoutType.BULLET_POINTS.value:
            self._create_bullet_points_slide(slide_data)
        elif slide_type == SlideLayoutType.TWO_COLUMN.value:
            self._create_two_column_slide(slide_data)
        elif slide_type == SlideLayoutType.CONTENT_WITH_IMAGE.value:
            self._create_content_with_image_slide(slide_data)
        else:
            # Default to bullet points for unknown types
            self._create_bullet_points_slide(slide_data)
    
    def _create_title_slide(self, slide_data: Dict[str, Any]):
        """Create a title slide"""
        
        # Use title slide layout
        slide_layout = self.presentation.slide_layouts[PPTXConstants.TITLE_LAYOUT_INDEX]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Apply background color if specified
        self._apply_background(slide)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data["title"]
        self._format_title(title, is_title_slide=True)
        
        # Set subtitle with content
        if slide.placeholders[1]:
            subtitle = slide.placeholders[1]
            subtitle.text = "\n".join(slide_data.get("content", []))
            self._format_subtitle(subtitle)
        
        # Add reference if it exists
        if "reference" in slide_data and slide_data["reference"]:
            self._add_reference(slide, slide_data["reference"])
        
        # Add speaker notes if they exist
        if "notes" in slide_data and slide_data["notes"]:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]
    
    def _create_bullet_points_slide(self, slide_data: Dict[str, Any]):
        """Create a slide with bullet points"""
        
        # Use content slide layout
        slide_layout = self.presentation.slide_layouts[PPTXConstants.CONTENT_LAYOUT_INDEX]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Apply background color if specified
        self._apply_background(slide)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data["title"]
        self._format_title(title)
        
        # Set content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, bullet_point in enumerate(slide_data.get("content", [])):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet_point
            p.level = 0
            self._format_bullet_point(p)
        
        # Add reference if it exists
        if "reference" in slide_data and slide_data["reference"]:
            self._add_reference(slide, slide_data["reference"])
        
        # Add speaker notes if they exist
        if "notes" in slide_data and slide_data["notes"]:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]
    
    def _create_two_column_slide(self, slide_data: Dict[str, Any]):
        """Create a slide with two columns"""
        
        # Use comparison layout which has two content placeholders
        slide_layout = self.presentation.slide_layouts[PPTXConstants.TWO_COLUMN_LAYOUT_INDEX]  # Comparison layout (usually)
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Apply background color if specified
        self._apply_background(slide)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data["title"]
        self._format_title(title)
        
        # Get content - split it for the two columns
        all_content = slide_data.get("content", [])
        mid_point = len(all_content) // 2
        
        left_content = all_content[:mid_point]
        right_content = all_content[mid_point:]
        
        # Try to find the two content placeholders
        content_placeholders = [shape for shape in slide.placeholders 
                               if shape.placeholder_format.idx in (1, 2)]
        
        if len(content_placeholders) >= 2:
            # Left column
            left_placeholder = content_placeholders[0]
            text_frame = left_placeholder.text_frame
            text_frame.clear()
            
            for i, bullet_point in enumerate(left_content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet_point
                p.level = 0
                self._format_bullet_point(p)
            
            # Right column
            right_placeholder = content_placeholders[1]
            text_frame = right_placeholder.text_frame
            text_frame.clear()
            
            for i, bullet_point in enumerate(right_content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet_point
                p.level = 0
                self._format_bullet_point(p)
        else:
            # Fallback if the expected layout is not available
            self._create_bullet_points_slide(slide_data)
        
        # Add reference if it exists
        if "reference" in slide_data and slide_data["reference"]:
            self._add_reference(slide, slide_data["reference"])
        
        # Add speaker notes if they exist
        if "notes" in slide_data and slide_data["notes"]:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]
    
    def _create_content_with_image_slide(self, slide_data: Dict[str, Any]):
        """Create a slide with content and an image placeholder"""
        
        # Use picture with caption layout
        slide_layout = self.presentation.slide_layouts[PPTXConstants.PICTURE_LAYOUT_INDEX]  # Usually picture with caption
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Apply background color if specified
        self._apply_background(slide)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data["title"]
        self._format_title(title)
        
        # Try to find content and picture placeholders
        placeholders = slide.placeholders
        
        # Add content to text placeholder if available
        text_placeholder = None
        for shape in placeholders:
            if shape.placeholder_format.type == 2:  # Content placeholder
                text_placeholder = shape
                break
        
        if text_placeholder:
            text_frame = text_placeholder.text_frame
            text_frame.clear()
            
            for i, bullet_point in enumerate(slide_data.get("content", [])):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet_point
                p.level = 0
                self._format_bullet_point(p)
        else:
            # Fallback if the expected layout is not available
            self._create_bullet_points_slide(slide_data)
        
        # Add reference if it exists
        if "reference" in slide_data and slide_data["reference"]:
            self._add_reference(slide, slide_data["reference"])
        
        # Add speaker notes if they exist
        if "notes" in slide_data and slide_data["notes"]:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]
    
    def _apply_background(self, slide):
        """Apply background color to slide if specified in config"""
        if 'background_color' in self.config:
            bg_color = self.config['background_color'].lstrip('#')
            if len(bg_color) == 6:
                try:
                    r = int(bg_color[0:2], 16)
                    g = int(bg_color[2:4], 16)
                    b = int(bg_color[4:6], 16)
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
                except ValueError:
                    # If there's a problem with the hex color, just skip it
                    pass
    
    def _add_reference(self, slide, reference_text):
        """Add a reference to the bottom of the slide"""
        left = Inches(PPTXConstants.REFERENCE_LEFT)
        top = Inches(PPTXConstants.REFERENCE_TOP)  # Position at bottom of slide
        width = Inches(PPTXConstants.REFERENCE_WIDTH)
        height = Inches(PPTXConstants.REFERENCE_HEIGHT)
        
        # Add a text box for the reference
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        # Add and format the reference text
        p = text_frame.paragraphs[0]
        p.text = reference_text
        p.alignment = PP_ALIGN.RIGHT  # Right-align the reference
        
        # Format the text to be smaller and lighter
        run = p.runs[0]
        run.font.size = Pt(PPTXConstants.REFERENCE_FONT_SIZE)
        run.font.name = self._get_font_name()
        
        # Use a lighter color for the reference text
        colors = self._get_theme_colors()
        if 'text' in colors:
            # Make it slightly lighter than normal text
            text_color = colors['text']
            # RGBColor is a tuple, so access values by index
            r = min(255, text_color[0] + 40)
            g = min(255, text_color[1] + 40)
            b = min(255, text_color[2] + 40)
            run.font.color.rgb = RGBColor(r, g, b)
        else:
            run.font.color.rgb = colors['secondary']
    
    def _get_theme_colors(self):
        """Get theme colors based on config or defaults"""
        theme_name = self.config.get('theme', PresentationTheme.PROFESSIONAL.value)
        
        # Default professional colors
        colors = {
            'primary': RGBColor(31, 73, 125),     # Dark blue
            'secondary': RGBColor(68, 114, 196),  # Medium blue
            'accent': RGBColor(237, 125, 49),     # Orange
            'text': RGBColor(68, 84, 106),        # Dark gray-blue
            'background': RGBColor(255, 255, 255) # White
        }
        
        # Apply theme-specific colors first
        if isinstance(theme_name, str):
            theme_lower = theme_name.lower()
            if theme_lower == PresentationTheme.CORPORATE.value:
                colors['primary'] = RGBColor(0, 50, 98)      # Dark blue
                colors['secondary'] = RGBColor(0, 118, 189)  # Medium blue
                colors['accent'] = RGBColor(242, 80, 34)     # Red/orange
            elif theme_lower == PresentationTheme.CREATIVE.value:
                colors['primary'] = RGBColor(185, 9, 11)     # Dark red
                colors['secondary'] = RGBColor(247, 150, 70) # Orange
                colors['accent'] = RGBColor(75, 172, 198)    # Light blue
            elif theme_lower == PresentationTheme.ACADEMIC.value:
                colors['primary'] = RGBColor(80, 16, 22)     # Maroon
                colors['secondary'] = RGBColor(155, 155, 155) # Gray
                colors['accent'] = RGBColor(200, 178, 115)   # Gold
            elif theme_lower == PresentationTheme.DARK.value or 'dark mode' in theme_lower:
                colors['primary'] = RGBColor(45, 45, 45)     # Dark gray
                colors['secondary'] = RGBColor(210, 210, 210) # Light gray
                colors['accent'] = RGBColor(52, 152, 219)    # Blue
                colors['text'] = RGBColor(240, 240, 240)     # Almost white
                colors['background'] = RGBColor(25, 25, 25)  # Almost black
            elif theme_lower == PresentationTheme.MINIMAL.value:
                colors['primary'] = RGBColor(40, 40, 40)     # Dark gray
                colors['secondary'] = RGBColor(120, 120, 120) # Medium gray
                colors['accent'] = RGBColor(200, 200, 200)   # Light gray
            elif PresentationTheme.MODERN_STARTUP.value in theme_lower:
                colors['primary'] = RGBColor(41, 128, 185)   # Bright blue
                colors['secondary'] = RGBColor(52, 73, 94)   # Dark blue-gray
                colors['accent'] = RGBColor(243, 156, 18)    # Orange
            elif PresentationTheme.YOUTHFUL.value in theme_lower:
                colors['primary'] = RGBColor(155, 89, 182)   # Purple
                colors['secondary'] = RGBColor(52, 152, 219) # Blue
                colors['accent'] = RGBColor(46, 204, 113)    # Green
        
        # Override with custom colors if specified
        if 'background_color' in self.config:
            bg_color = self.config['background_color'].lstrip('#')
            if len(bg_color) == 6:
                try:
                    r = int(bg_color[0:2], 16)
                    g = int(bg_color[2:4], 16)
                    b = int(bg_color[4:6], 16)
                    colors['background'] = RGBColor(r, g, b)
                    
                    # If background is dark, use light text by default
                    brightness = (r * 299 + g * 587 + b * 114) / 1000
                    if brightness < 128:  # Dark background
                        colors['text'] = RGBColor(240, 240, 240)  # Light text
                except ValueError:
                    pass
        
        # Apply custom title color (maps to primary)
        if 'title_color' in self.config:
            title_color = self.config['title_color'].lstrip('#')
            if len(title_color) == 6:
                try:
                    r = int(title_color[0:2], 16)
                    g = int(title_color[2:4], 16)
                    b = int(title_color[4:6], 16)
                    colors['primary'] = RGBColor(r, g, b)
                except ValueError:
                    pass
        
        # Apply custom content color (maps to text)
        if 'content_color' in self.config:
            content_color = self.config['content_color'].lstrip('#')
            if len(content_color) == 6:
                try:
                    r = int(content_color[0:2], 16)
                    g = int(content_color[2:4], 16)
                    b = int(content_color[4:6], 16)
                    colors['text'] = RGBColor(r, g, b)
                    colors['secondary'] = RGBColor(r, g, b)  # Also update secondary
                except ValueError:
                    pass
        
        # Apply custom accent color
        if 'accent_color' in self.config:
            accent_color = self.config['accent_color'].lstrip('#')
            if len(accent_color) == 6:
                try:
                    r = int(accent_color[0:2], 16)
                    g = int(accent_color[2:4], 16)
                    b = int(accent_color[4:6], 16)
                    colors['accent'] = RGBColor(r, g, b)
                except ValueError:
                    pass
        
        return colors
    
    def _get_font_name(self):
        """Get font name from config or default"""
        return self.config.get('font', self.config.get('font_name', PPTXConstants.DEFAULT_FONT))
    
    def _format_title(self, title_shape, is_title_slide=False):
        """Format title text based on theme and config"""
        if title_shape.has_text_frame:
            colors = self._get_theme_colors()
            font_name = self._get_font_name()
            
            text_frame = title_shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(PPTXConstants.TITLE_SLIDE_FONT_SIZE if is_title_slide else PPTXConstants.TITLE_FONT_SIZE)
                    run.font.bold = True
                    run.font.color.rgb = colors['primary']
    
    def _format_subtitle(self, subtitle_shape):
        """Format subtitle text based on theme and config"""
        if subtitle_shape.has_text_frame:
            colors = self._get_theme_colors()
            font_name = self._get_font_name()
            
            text_frame = subtitle_shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(PPTXConstants.SUBTITLE_FONT_SIZE)
                    run.font.color.rgb = colors['secondary']
    
    def _format_bullet_point(self, paragraph, is_conclusion=False):
        """Format bullet point text based on theme and config"""
        colors = self._get_theme_colors()
        font_name = self._get_font_name()
        
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(PPTXConstants.CONCLUSION_FONT_SIZE if is_conclusion else PPTXConstants.BULLET_FONT_SIZE)
            run.font.color.rgb = colors['text'] if 'text' in colors else colors['secondary']
            if is_conclusion:
                run.font.bold = True 